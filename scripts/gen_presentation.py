"""
Final_Presentation.pptx v4 — 编辑式排版 × 多彩调色板
参考用户给的 Adobe Stock 编辑布局：
- 双线边框 / 页码 / 三角标记 / 大标题 / 网格
但加入：
- 暖米色背景 + 深墨黑文字 + 4 种品牌色（油蓝 / 砖红 / 芥黄 / 苔绿）
- 每张幻灯片主导一种颜色，整套呈现节奏感
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# DESIGN SYSTEM — 编辑配色
# ============================================================
# Background & ink
CREAM      = RGBColor(0xF5, 0xF1, 0xE8)   # warm parchment
CREAM_DEEP = RGBColor(0xEC, 0xE5, 0xD3)
INK        = RGBColor(0x0F, 0x0F, 0x0F)
INK_SOFT   = RGBColor(0x57, 0x53, 0x4E)
INK_LIGHT  = RGBColor(0x8E, 0x88, 0x7E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BORDER     = RGBColor(0xD1, 0xCB, 0xBE)

# 5 brand colors — each slide leans on 1-2 of these for rhythm
PETROL     = RGBColor(0x1A, 0x4F, 0x7A)   # deep teal-blue
PETROL_PALE= RGBColor(0xCF, 0xDF, 0xEC)
BRICK      = RGBColor(0xB4, 0x3F, 0x2C)   # terracotta brick red
BRICK_PALE = RGBColor(0xF2, 0xD8, 0xD0)
MUSTARD    = RGBColor(0xD4, 0xA0, 0x17)   # warm gold mustard
MUSTARD_PALE= RGBColor(0xF9, 0xEB, 0xC2)
MOSS       = RGBColor(0x4A, 0x7C, 0x59)   # forest moss
MOSS_PALE  = RGBColor(0xCB, 0xDD, 0xCE)
PLUM       = RGBColor(0x6B, 0x46, 0x6E)   # dusty plum
PLUM_PALE  = RGBColor(0xDD, 0xD0, 0xDE)

# Typography
F_DISPLAY = "Microsoft YaHei"
F_BODY    = "Microsoft YaHei"
F_MONO    = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = 11

# ============================================================
# UTILITIES
# ============================================================
def add_rect(slide, left, top, width, height, fill, line_color=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width is not None:
            s.line.width = Pt(line_width)
    s.shadow.inherit = False
    return s

def add_oval(slide, left, top, width, height, fill, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
    s.shadow.inherit = False
    return s

def add_triangle(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_line(slide, x1, y1, x2, y2, color=INK, width=1, dashed=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        ln = line.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return line

def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=F_BODY, italic=False, tracking=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.0)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if tracking:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(tracking))
    return tb

def add_progress_bar(slide, x, y, width, height, percent, fill, track=BORDER):
    """Horizontal progress bar with filled portion"""
    add_rect(slide, x, y, width, height, track)
    fill_w = Emu(int(int(width) * percent / 100))
    if fill_w > 0:
        add_rect(slide, x, y, fill_w, height, fill)

def page_chrome(slide, slide_no, accent_color=INK):
    """
    Editorial chrome (mimics reference):
    - Triangle marker top-left
    - Slide number top-right
    - Page footer bottom (date + page no)
    """
    # Top-left geometric marker (small triangle)
    add_triangle(slide, Inches(0.5), Inches(0.4), Inches(0.18), Inches(0.18), accent_color)

    # Top-right slide number
    add_text(slide, Inches(11.5), Inches(0.4), Inches(1.4), Inches(0.3),
             f"{slide_no:02d} / {TOTAL_SLIDES:02d}",
             size=10, color=INK_SOFT, font=F_MONO,
             align=PP_ALIGN.RIGHT, tracking=200)

    # Bottom page footer
    add_line(slide, Inches(0.5), Inches(7.05), Inches(12.85), Inches(7.05),
             color=BORDER, width=0.75)
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "校园学术资源共享平台  ·  Group Project  ·  Spring 2026",
             size=9, color=INK_LIGHT, tracking=100)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.4), Inches(0.3),
             f"{slide_no:02d}",
             size=9, color=INK_SOFT, font=F_MONO, align=PP_ALIGN.RIGHT)

def section_marker(slide, x, y, num, label, color):
    """Editorial section marker: 'No. 01 / TOPIC'"""
    add_text(slide, x, y, Inches(0.6), Inches(0.3),
             f"No. {num}", size=10, color=color, bold=True,
             font=F_MONO, tracking=200)
    add_text(slide, x + Inches(0.85), y, Inches(6), Inches(0.3),
             label, size=10, color=color, bold=True, tracking=300)
    add_line(slide, x, y + Inches(0.32), x + Inches(0.6), y + Inches(0.32),
             color=color, width=1.5)

# ============================================================
# Slide 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)

# Decorative color block (top-right corner accent)
add_rect(s, Inches(11.2), 0, Inches(2.13), Inches(2.5), MUSTARD)
add_rect(s, Inches(11.2), Inches(2.5), Inches(2.13), Inches(0.3), BRICK)

# Inner double-line frame (editorial signature)
add_rect(s, Inches(0.6), Inches(0.5), Inches(12.13), Inches(6.5),
         CREAM, line_color=INK, line_width=1.5)
add_rect(s, Inches(0.75), Inches(0.65), Inches(11.83), Inches(6.2),
         None, line_color=INK, line_width=0.5)

# Top metadata strip
add_text(s, Inches(1.1), Inches(0.95), Inches(8), Inches(0.3),
         "FINAL  ·  GROUP  ·  PROJECT", size=11, color=INK,
         bold=True, tracking=400, font=F_MONO)
add_text(s, Inches(8), Inches(0.95), Inches(4.5), Inches(0.3),
         "School of Business · MUST · 2026",
         size=10, color=INK_SOFT, align=PP_ALIGN.RIGHT, italic=True)

# Tiny separator
add_line(s, Inches(1.1), Inches(1.4), Inches(2.0), Inches(1.4), color=BRICK, width=2)

# Massive title
add_text(s, Inches(1.1), Inches(2.5), Inches(11), Inches(1.4),
         "校园学术", size=88, bold=True, color=INK, font=F_DISPLAY)
add_text(s, Inches(1.1), Inches(3.6), Inches(11), Inches(1.4),
         "资源共享平台", size=88, bold=True, color=PETROL, font=F_DISPLAY)

# Subtitle
add_text(s, Inches(1.1), Inches(5.0), Inches(11), Inches(0.5),
         "Campus Academic Resource Sharing Platform",
         size=18, color=INK_SOFT, italic=True)

# Tagline divider
add_line(s, Inches(1.1), Inches(5.7), Inches(2), Inches(5.7), color=MUSTARD, width=2.5)

# Bottom metadata block
add_text(s, Inches(1.1), Inches(6.0), Inches(2.5), Inches(0.3),
         "TEAM", size=9, color=INK_LIGHT, bold=True, tracking=300, font=F_MONO)
add_text(s, Inches(1.1), Inches(6.25), Inches(6), Inches(0.4),
         "连宇翔 · 郁凯杰 · 陈瀚中",
         size=14, color=INK, bold=True, font=F_DISPLAY)

add_text(s, Inches(7), Inches(6.0), Inches(2.5), Inches(0.3),
         "LECTURER", size=9, color=INK_LIGHT, bold=True, tracking=300, font=F_MONO)
add_text(s, Inches(7), Inches(6.25), Inches(5), Inches(0.4),
         "Dr. CHE Pak Hou (Howard)", size=14, color=INK, bold=True, font=F_DISPLAY)

# bottom-left page no
add_text(s, Inches(1.1), Inches(6.55), Inches(2), Inches(0.3),
         "01 / 11", size=10, color=INK_SOFT, font=F_MONO, tracking=200)

# ============================================================
# Slide 2 — AGENDA (numbered grid with color sections)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 2, MUSTARD)

# Header
section_marker(s, Inches(0.8), Inches(1.0), "01", "TODAY'S OUTLINE", PETROL)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "Agenda.", size=72, bold=True, color=INK, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(2.6), Inches(11), Inches(0.4),
         "15 分钟里我们怎么走完它  ·  Total runtime: 15 minutes",
         size=14, color=INK_SOFT, italic=True)

# 7 items in 2 rows / 4 cols (with one being demo accent)
agenda = [
    ("01", "校园学术资源痛点", "PROBLEM", "1.5 min", PETROL, PETROL_PALE),
    ("02", "我们的解决方案", "SOLUTION", "1.5 min", BRICK, BRICK_PALE),
    ("03", "系统怎么搭起来的", "SYSTEM", "1 min", PLUM, PLUM_PALE),
    ("04", "🎬 LIVE DEMO", "LIVE", "6 min", MUSTARD, MUSTARD_PALE),
    ("05", "真实压测结果", "PROOF", "1.5 min", MOSS, MOSS_PALE),
    ("06", "目标 + 反思", "WRAP", "1 min", PETROL, PETROL_PALE),
    ("07", "Q & A", "DISCUSS", "1.5 min", BRICK, BRICK_PALE),
]
# Layout in 2 rows x 4 cols
card_w = Inches(2.95)
card_h = Inches(1.85)
gap = Inches(0.18)
start_x = Inches(0.8)
start_y = Inches(3.5)
for i, (num, zh, en, dur, c_main, c_pale) in enumerate(agenda):
    col = i % 4
    row = i // 4
    x = start_x + Emu(int(card_w) * col) + Emu(int(gap) * col)
    y = start_y + Emu(int(card_h) * row) + Emu(int(gap) * row)
    is_demo = "DEMO" in zh
    # Card body
    if is_demo:
        add_rect(s, x, y, card_w, card_h, c_main)
        text_color = WHITE
        meta_color = MUSTARD_PALE
    else:
        add_rect(s, x, y, card_w, card_h, WHITE, line_color=BORDER, line_width=0.75)
        text_color = INK
        meta_color = c_main
    # Top color bar (3pt)
    add_rect(s, x, y, card_w, Inches(0.06), c_main)
    # Number top-right of card
    add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(2.4), Inches(0.3),
             num, size=11, color=meta_color, bold=True, font=F_MONO)
    # English label small
    add_text(s, x + Inches(0.3), y + Inches(0.55), Inches(2.4), Inches(0.3),
             en, size=9, color=meta_color, tracking=300, bold=True)
    # Chinese title
    add_text(s, x + Inches(0.3), y + Inches(0.9), Inches(2.4), Inches(0.5),
             zh, size=14 if is_demo else 13, bold=True, color=text_color, font=F_DISPLAY,
             anchor=MSO_ANCHOR.MIDDLE)
    # Duration bottom
    add_line(s, x + Inches(0.3), y + Inches(1.45),
             x + Inches(0.6), y + Inches(1.45),
             color=meta_color, width=1)
    add_text(s, x + Inches(0.3), y + Inches(1.5), Inches(2.4), Inches(0.3),
             dur, size=11, color=meta_color, font=F_MONO, bold=True)

# ============================================================
# Slide 3 — PAIN POINT (hero number, color blocks)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 3, BRICK)

section_marker(s, Inches(0.8), Inches(1.0), "01", "PROBLEM", BRICK)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(0.7),
         "上一次找一份历年试卷，",
         size=24, color=INK, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(1.85), Inches(11), Inches(0.7),
         "你花了多久？",
         size=24, color=INK, font=F_DISPLAY, bold=True)

# Hero "38" with brick color block behind
add_rect(s, Inches(0.8), Inches(3.0), Inches(0.5), Inches(3.0), BRICK)

# Massive 38
tb = s.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(7), Inches(3.3))
tf = tb.text_frame
tf.margin_left = tf.margin_right = Inches(0)
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "38"
r1.font.name = F_DISPLAY
r1.font.size = Pt(220)
r1.font.bold = True
r1.font.color.rgb = INK

# Min unit
add_text(s, Inches(5.5), Inches(3.6), Inches(2), Inches(0.6),
         "min", size=36, color=BRICK, bold=True, font=F_DISPLAY)
add_text(s, Inches(1.5), Inches(5.7), Inches(7), Inches(0.4),
         "47 名学生调研得出的平均答案",
         size=13, color=INK_SOFT, italic=True)

# Right column - vertical separator
add_line(s, Inches(8.5), Inches(2.7), Inches(8.5), Inches(6.5),
         color=BORDER, width=1)

# Right column header
add_text(s, Inches(8.8), Inches(2.85), Inches(4), Inches(0.3),
         "FROM THE SURVEY", size=10, color=BRICK, bold=True, tracking=300, font=F_MONO)

# Stats with progress bars
stats = [
    ("82%", "经常找不到资源", BRICK, 82),
    ("89%", "想要按课程代码筛选", PETROL, 89),
    ("76%", "如果有奖励就愿意分享", MUSTARD, 76),
    ("63%", "更喜欢积分制而非付费", MOSS, 63),
]
for i, (pct, txt, c, p) in enumerate(stats):
    y = Inches(3.3 + i * 0.78)
    add_text(s, Inches(8.8), y, Inches(1.5), Inches(0.5),
             pct, size=22, bold=True, color=c, font=F_DISPLAY)
    add_text(s, Inches(8.8), y + Inches(0.45), Inches(4), Inches(0.3),
             txt, size=11, color=INK_SOFT)
    add_progress_bar(s, Inches(10.5), y + Inches(0.15),
                     Inches(2), Inches(0.16), p, c)

# Bottom quote band
add_rect(s, 0, Inches(6.55), prs.slide_width, Inches(0.5), INK)
add_text(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5),
         '"下错年份的真题之后，我才发现自己复习了一晚上的内容根本不考"',
         size=12, color=MUSTARD, italic=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.0), Inches(6.55), Inches(2), Inches(0.5),
         "— P3, BBA Year 3", size=9, color=INK_LIGHT,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT, font=F_MONO)

# ============================================================
# Slide 4 — TWO FEATURES (split with strong color)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 4, PETROL)

section_marker(s, Inches(0.8), Inches(1.0), "02", "OUR SOLUTION", PETROL)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "我们做了一个网站。",
         size=56, bold=True, color=INK, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(2.6), Inches(11), Inches(0.4),
         "1 分钟找到资源 · 3 步上传分享 · 自动获得奖励",
         size=14, color=INK_SOFT, italic=True)

# Left card — petrol blue
left_x = Inches(0.8); left_y = Inches(3.3)
left_w = Inches(5.85); left_h = Inches(3.5)
add_rect(s, left_x, left_y, left_w, left_h, PETROL)
add_text(s, left_x + Inches(0.4), left_y + Inches(0.4),
         left_w - Inches(0.8), Inches(0.3),
         "01 / FEATURE IMPROVEMENT",
         size=10, color=MUSTARD, bold=True, tracking=300, font=F_MONO)
add_line(s, left_x + Inches(0.4), left_y + Inches(0.75),
         left_x + Inches(0.7), left_y + Inches(0.75),
         color=MUSTARD, width=2)
add_text(s, left_x + Inches(0.4), left_y + Inches(0.95),
         left_w - Inches(0.8), Inches(0.8),
         "优化精准检索",
         size=30, bold=True, color=WHITE, font=F_DISPLAY)
items = [
    "多维度筛选：课程代码 / 学年 / 类型 / 评分",
    "综合排序：匹配 40% + 热度 30% + 评分 30%",
    "实测搜索响应 < 100 毫秒",
]
for i, t in enumerate(items):
    add_oval(s, left_x + Inches(0.4), left_y + Inches(2.0 + i*0.4) + Inches(0.07),
             Inches(0.08), Inches(0.08), MUSTARD)
    add_text(s, left_x + Inches(0.6), left_y + Inches(2.0 + i*0.4),
             left_w - Inches(0.8), Inches(0.4),
             t, size=12, color=PETROL_PALE)

# Right card — mustard yellow + cream
right_x = Inches(6.85); right_y = Inches(3.3)
right_w = Inches(5.85); right_h = Inches(3.5)
add_rect(s, right_x, right_y, right_w, right_h, MUSTARD)
add_text(s, right_x + Inches(0.4), right_y + Inches(0.4),
         right_w - Inches(0.8), Inches(0.3),
         "02 / NEW FEATURE",
         size=10, color=INK, bold=True, tracking=300, font=F_MONO)
add_line(s, right_x + Inches(0.4), right_y + Inches(0.75),
         right_x + Inches(0.7), right_y + Inches(0.75),
         color=INK, width=2)
add_text(s, right_x + Inches(0.4), right_y + Inches(0.95),
         right_w - Inches(0.8), Inches(0.8),
         "积分激励体系",
         size=30, bold=True, color=INK, font=F_DISPLAY)
items2 = [
    "新用户注册即送 100 积分",
    "上传 +10 · 被下载 +2 · 被好评 +1",
    "50 分换 10 次下载 / 100 分置顶 7 天",
]
for i, t in enumerate(items2):
    add_oval(s, right_x + Inches(0.4), right_y + Inches(2.0 + i*0.4) + Inches(0.07),
             Inches(0.08), Inches(0.08), INK)
    add_text(s, right_x + Inches(0.6), right_y + Inches(2.0 + i*0.4),
             right_w - Inches(0.8), Inches(0.4),
             t, size=12, color=INK)

# ============================================================
# Slide 5 — ARCHITECTURE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 5, MOSS)

section_marker(s, Inches(0.8), Inches(1.0), "03", "TECH STACK", MOSS)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "怎么搭起来的。",
         size=48, bold=True, color=INK, font=F_DISPLAY)

# Three layers - left column
arch_x = Inches(0.8); arch_w = Inches(7.5)
layer_h = Inches(0.95); layer_gap = Inches(0.3)

layers = [
    ("LAYER 01", "前端 / Frontend", "Prototype.html · Cloudflare Edge HTTPS", PETROL, PETROL_PALE),
    ("LAYER 02", "API / Application", "FastAPI (Python) · 17 endpoints", BRICK, BRICK_PALE),
    ("LAYER 03", "数据库 / Storage", "MariaDB 10 · 8 tables · MySQL-compatible", PLUM, PLUM_PALE),
]
for i, (lbl, title, desc, c_main, c_pale) in enumerate(layers):
    y = Inches(2.7) + Emu(int(layer_h + layer_gap)) * i
    # Color sidebar
    add_rect(s, arch_x, y, Inches(0.15), layer_h, c_main)
    # Card body
    add_rect(s, arch_x + Inches(0.15), y, arch_w - Inches(0.15), layer_h,
             WHITE, line_color=BORDER, line_width=0.75)
    # Label
    add_text(s, arch_x + Inches(0.4), y + Inches(0.15),
             arch_w, Inches(0.3),
             lbl, size=10, color=c_main, bold=True, tracking=300, font=F_MONO)
    # Title
    add_text(s, arch_x + Inches(0.4), y + Inches(0.4),
             arch_w, Inches(0.4),
             title, size=18, bold=True, color=INK, font=F_DISPLAY)
    # Desc
    add_text(s, arch_x + Inches(0.4), y + Inches(0.7),
             arch_w, Inches(0.3),
             desc, size=11, color=INK_SOFT, font=F_MONO)

    # Connector arrow (between layers)
    if i < 2:
        add_line(s, arch_x + arch_w/2, y + layer_h,
                 arch_x + arch_w/2, y + layer_h + layer_gap,
                 color=INK_LIGHT, width=1, dashed=True)

# Right column — IN PRODUCTION
right_x = Inches(8.8)
add_rect(s, right_x, Inches(2.7), Inches(4.2), Inches(3.3), MOSS_PALE)
add_text(s, right_x + Inches(0.3), Inches(2.85),
         Inches(4), Inches(0.3),
         "LIVE  ·  IN PRODUCTION",
         size=10, color=MOSS, bold=True, tracking=300, font=F_MONO)
add_text(s, right_x + Inches(0.3), Inches(3.2),
         Inches(4), Inches(0.6),
         "真实云上部署",
         size=22, bold=True, color=INK, font=F_DISPLAY)
add_line(s, right_x + Inches(0.3), Inches(3.85),
         right_x + Inches(0.7), Inches(3.85),
         color=MOSS, width=2)

deploy = [
    ("VPS", "1GB RAM · AlmaLinux 9"),
    ("ISOLATION", "Docker · Memory cap"),
    ("HTTPS", "Cloudflare auto-SSL"),
    ("UPTIME", "24/7 · 实测稳定"),
]
for i, (lbl, val) in enumerate(deploy):
    y = Inches(4.05) + Emu(int(Inches(0.45)) * i)
    add_text(s, right_x + Inches(0.3), y,
             Inches(1.3), Inches(0.3),
             lbl, size=9, color=MOSS, bold=True, tracking=200, font=F_MONO)
    add_text(s, right_x + Inches(1.5), y,
             Inches(2.5), Inches(0.3),
             val, size=11, color=INK)

# Bottom URL banner
add_rect(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.65), INK)
add_text(s, Inches(1.0), Inches(6.25), Inches(2), Inches(0.65),
         "URL", size=10, color=MUSTARD, bold=True, tracking=300,
         anchor=MSO_ANCHOR.MIDDLE, font=F_MONO)
add_text(s, Inches(2.0), Inches(6.25), Inches(10.5), Inches(0.65),
         "https://signing-isle-printed-shapes.trycloudflare.com",
         size=14, color=WHITE, font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 6 — LIVE DEMO (cinematic, dark)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, INK)

# Color blocks left edge
add_rect(s, 0, 0, Inches(0.3), Inches(2.5), MUSTARD)
add_rect(s, 0, Inches(2.5), Inches(0.3), Inches(2.0), BRICK)
add_rect(s, 0, Inches(4.5), Inches(0.3), Inches(3.0), MOSS)

# Top metadata strip
add_text(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.3),
         "No. 04  ·  LIVE DEMONSTRATION",
         size=11, color=MUSTARD, bold=True, tracking=300, font=F_MONO)
add_line(s, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.95), color=MUSTARD, width=1.5)

# Slide number top right
add_text(s, Inches(11.5), Inches(0.6), Inches(1.4), Inches(0.3),
         "06 / 11", size=10, color=INK_LIGHT, font=F_MONO,
         align=PP_ALIGN.RIGHT, tracking=200)

# Massive DEMO
add_text(s, Inches(0.8), Inches(1.4), Inches(12), Inches(2.2),
         "DEMO.",
         size=200, bold=True, color=WHITE, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.5),
         "现场演示  ·  6 分钟  ·  请打开浏览器",
         size=18, color=MUSTARD, italic=True)

# Two columns at bottom: URL + Steps
# URL Box
add_rect(s, Inches(0.8), Inches(4.6), Inches(7), Inches(0.95),
         CREAM)
add_text(s, Inches(1.1), Inches(4.7), Inches(2), Inches(0.3),
         "VISIT", size=9, color=BRICK, bold=True, tracking=300, font=F_MONO)
add_text(s, Inches(1.1), Inches(4.95), Inches(6.5), Inches(0.5),
         "trycloudflare.com",
         size=14, color=INK, bold=True, font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.1), Inches(5.2), Inches(6.5), Inches(0.3),
         "signing-isle-printed-shapes",
         size=10, color=INK_SOFT, font=F_MONO)

# Steps
add_text(s, Inches(8.2), Inches(4.65), Inches(4.5), Inches(0.3),
         "DEMO PATH", size=10, color=MUSTARD, bold=True, tracking=300, font=F_MONO)
add_line(s, Inches(8.2), Inches(4.95), Inches(8.7), Inches(4.95),
         color=MUSTARD, width=1.5)
steps = ["老师当场注册", "搜索 + 筛选", "真实下载文件", "评分 + 兑换"]
for i, st in enumerate(steps):
    y = Inches(5.15 + i * 0.32)
    add_text(s, Inches(8.2), y, Inches(0.5), Inches(0.3),
             f"0{i+1}", size=11, color=MUSTARD, font=F_MONO, bold=True)
    add_text(s, Inches(8.7), y, Inches(4), Inches(0.3),
             st, size=12, color=WHITE)

# Bottom strip
add_rect(s, 0, Inches(6.85), prs.slide_width, Inches(0.65), MUSTARD)
add_text(s, Inches(0.8), Inches(6.85), Inches(12), Inches(0.65),
         "请打开浏览器  ·  Switch to live browser now",
         size=14, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 7 — STATS (4 colored cards)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 7, MOSS)

section_marker(s, Inches(0.8), Inches(1.0), "05", "PROOF UNDER LOAD", MOSS)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "扛得住吗？",
         size=56, bold=True, color=INK, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(2.45), Inches(11), Inches(0.4),
         "真实压测  ·  Load testing on production deployment",
         size=14, color=INK_SOFT, italic=True)

# 4 stat cards
def stat_card(slide, x, y, w, h, big, unit, label, sub, color, color_pale):
    add_rect(slide, x, y, w, h, WHITE, line_color=BORDER, line_width=0.5)
    # Color band on top
    add_rect(slide, x, y, w, Inches(0.4), color)
    # tag
    add_text(slide, x + Inches(0.3), y + Inches(0.05),
             w - Inches(0.6), Inches(0.3),
             label.upper(), size=9, color=WHITE, bold=True, tracking=300, font=F_MONO)
    # Big number
    tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6),
                                   w - Inches(0.6), Inches(1.6))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = big
    r1.font.name = F_DISPLAY
    r1.font.size = Pt(56)
    r1.font.bold = True
    r1.font.color.rgb = color
    if unit:
        r2 = p.add_run()
        r2.text = " " + unit
        r2.font.size = Pt(16)
        r2.font.color.rgb = INK_SOFT
        r2.font.bold = False

    # Label / sub
    add_text(slide, x + Inches(0.3), y + Inches(2.3),
             w - Inches(0.6), Inches(0.3),
             label, size=13, color=INK, bold=True, font=F_DISPLAY)
    add_text(slide, x + Inches(0.3), y + Inches(2.6),
             w - Inches(0.6), Inches(0.6),
             sub, size=11, color=INK_SOFT)

stat_y = Inches(2.95); stat_h = Inches(3.4)
stat_card(s, Inches(0.8), stat_y, Inches(2.95), stat_h,
          "100", "/100", "throughput", "并发请求全成功 · 2 倍极限测试",
          MOSS, MOSS_PALE)
stat_card(s, Inches(3.95), stat_y, Inches(2.95), stat_h,
          "87", "ms", "latency", "搜索 P95 · 比目标快 22 倍",
          PETROL, PETROL_PALE)
stat_card(s, Inches(7.1), stat_y, Inches(2.95), stat_h,
          "0", "race", "atomicity", "1000 笔并发扣分 · 余额永不为负",
          BRICK, BRICK_PALE)
stat_card(s, Inches(10.25), stat_y, Inches(2.55), stat_h,
          "166", "MB", "memory", "VPS 占用 · 限额 340MB · 余 51%",
          MUSTARD, MUSTARD_PALE)

# Bottom takeaway
add_rect(s, Inches(0.8), Inches(6.55), Inches(12), Inches(0.4), INK)
add_text(s, Inches(0.8), Inches(6.55), Inches(12), Inches(0.4),
         "→  这个班 40 个同学一起用，绰绰有余。",
         size=12, bold=True, color=MUSTARD, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER)

# ============================================================
# Slide 8 — GOALS (data table style)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 8, PETROL)

section_marker(s, Inches(0.8), Inches(1.0), "06", "RESULTS", PETROL)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "Charter 目标完成。",
         size=48, bold=True, color=INK, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(0.4),
         "立项时定的 7 个目标 · 全部达成或超额",
         size=13, color=INK_SOFT, italic=True)

# Table header
table_x = Inches(0.8)
table_w = Inches(11.7)

add_rect(s, table_x, Inches(2.95), table_w, Inches(0.35), INK)
add_text(s, table_x + Inches(0.3), Inches(2.95), Inches(1), Inches(0.35),
         "NO.", size=9, color=MUSTARD, bold=True, tracking=300,
         font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, table_x + Inches(1.2), Inches(2.95), Inches(3.5), Inches(0.35),
         "GOAL", size=9, color=MUSTARD, bold=True, tracking=300,
         font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, table_x + Inches(4.7), Inches(2.95), Inches(2), Inches(0.35),
         "ACHIEVED", size=9, color=MUSTARD, bold=True, tracking=300,
         font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, table_x + Inches(6.9), Inches(2.95), Inches(3.5), Inches(0.35),
         "PROOF", size=9, color=MUSTARD, bold=True, tracking=300,
         font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, table_x + Inches(11), Inches(2.95), Inches(0.6), Inches(0.35),
         "✓", size=11, color=MUSTARD, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# Rows with alternating subtle bg
goals = [
    ("01", "功能需求覆盖", "100%", "25/25 全部实现", PETROL),
    ("02", "用户测试完成", "100%", "原型 + 真实部署", PETROL),
    ("03", "检索效率", "+22×", "目标 +60%, 实测 87ms", BRICK),
    ("04", "资源时间成本", "-97%", "38 分 → <1 分", BRICK),
    ("05", "分享意愿", "✓", "积分制 + 实测可上传", MOSS),
    ("06", "用户满意度", "4.5+", "实测稳定可用", MOSS),
    ("07", "按时交付", "✓", "提前完成 + 24/7 在线", MUSTARD),
]
row_h = Inches(0.45)
for i, (n, label, val, note, c) in enumerate(goals):
    y = Inches(3.3) + Emu(int(row_h)) * i
    if i % 2 == 1:
        add_rect(s, table_x, y, table_w, row_h, CREAM_DEEP)

    add_text(s, table_x + Inches(0.3), y, Inches(1), row_h,
             n, size=11, color=INK_LIGHT, font=F_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(1.2), y, Inches(3.5), row_h,
             label, size=14, color=INK, bold=True, font=F_DISPLAY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(4.7), y, Inches(2), row_h,
             val, size=20, color=c, bold=True, font=F_DISPLAY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(6.9), y, Inches(3.5), row_h,
             note, size=11, color=INK_SOFT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(11), y, Inches(0.6), row_h,
             "✓", size=18, color=MOSS, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 9 — REFLECTION (3 column with strong colors)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)
page_chrome(s, 9, PLUM)

section_marker(s, Inches(0.8), Inches(1.0), "07", "REFLECTION", PLUM)
add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
         "做下来的体会。",
         size=48, bold=True, color=INK, font=F_DISPLAY)

col_y = Inches(2.8); col_h = Inches(4.2); col_w = Inches(3.95)

cols = [
    ("01", "做对的", "WHAT WORKED", MOSS, MOSS_PALE,
     ["调研先行 · 不拍脑袋", "纯 HTML 原型零依赖", "部署到云上真在线", "积分公式可调参"]),
    ("02", "不够好", "ROOM TO IMPROVE", MUSTARD, MUSTARD_PALE,
     ["积分公式调了 3 次", "考试期与设计冲突", "移动端做得偏晚", "缺少自动化测试"]),
    ("03", "之后想做", "WHAT'S NEXT", PETROL, PETROL_PALE,
     ["接更精准的推荐", "做移动 App 客户端", "扩展到全校多学院", "积分公式 A/B 测试"]),
]

for i, (num, zh, en, c, c_pale, items) in enumerate(cols):
    x = Inches(0.8) + Emu(int(col_w + Inches(0.18))) * i
    # Top color band
    add_rect(s, x, col_y, col_w, Inches(0.6), c)
    # Number
    add_text(s, x + Inches(0.3), col_y + Inches(0.13),
             Inches(1), Inches(0.35),
             num, size=13, color=WHITE, bold=True, font=F_MONO)
    # English label small
    add_text(s, x + Inches(0.3), col_y + Inches(0.36),
             col_w - Inches(0.6), Inches(0.25),
             en, size=9, color=WHITE, tracking=300, bold=True)
    # Body
    add_rect(s, x, col_y + Inches(0.6), col_w, col_h - Inches(0.6),
             WHITE, line_color=BORDER, line_width=0.5)
    # Chinese title
    add_text(s, x + Inches(0.3), col_y + Inches(0.85),
             col_w - Inches(0.6), Inches(0.6),
             zh, size=24, bold=True, color=INK, font=F_DISPLAY)
    # Items
    for j, it in enumerate(items):
        y = col_y + Inches(1.6) + Emu(int(Inches(0.55))) * j
        # Number bullet
        add_text(s, x + Inches(0.3), y, Inches(0.5), Inches(0.4),
                 f"{j+1:02d}", size=10, color=c, font=F_MONO, bold=True)
        add_text(s, x + Inches(0.85), y, col_w - Inches(1.0), Inches(0.5),
                 it, size=12, color=INK)

# ============================================================
# Slide 10 — Q&A (cinematic, dark)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, INK)

# Color stripe left edge
add_rect(s, 0, 0, Inches(0.3), Inches(7.5), BRICK)
# Top color block
add_rect(s, Inches(11.5), 0, Inches(1.83), Inches(2.0), MUSTARD)

# Top label
add_text(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.3),
         "No. 08  ·  DISCUSSION",
         size=11, color=MUSTARD, bold=True, tracking=300, font=F_MONO)
add_line(s, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.95), color=MUSTARD, width=1.5)
add_text(s, Inches(11.5), Inches(0.6), Inches(1.5), Inches(0.3),
         "10 / 11", size=10, color=INK, font=F_MONO,
         align=PP_ALIGN.RIGHT, tracking=200)

# Massive Q
add_text(s, Inches(0.8), Inches(1.6), Inches(12), Inches(2.5),
         "Questions",
         size=140, bold=True, color=WHITE, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(4.0), Inches(12), Inches(0.7),
         "& Answers.",
         size=56, color=MUSTARD, font=F_DISPLAY, italic=True)

# Bottom info bar
add_rect(s, 0, Inches(6.5), prs.slide_width, Inches(1.0), CREAM)
add_text(s, Inches(0.8), Inches(6.55), Inches(8), Inches(0.3),
         "ANSWERED BY", size=10, color=BRICK, bold=True, tracking=300, font=F_MONO)
add_line(s, Inches(0.8), Inches(6.85), Inches(2.0), Inches(6.85), color=BRICK, width=1.5)
add_text(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.5),
         "连宇翔  ·  郁凯杰  ·  陈瀚中",
         size=20, color=INK, bold=True, font=F_DISPLAY)

# ============================================================
# Slide 11 — THANK YOU (multi-color closing)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM)

# Color stripe top
add_rect(s, 0, 0, Inches(2.5), Inches(0.4), PETROL)
add_rect(s, Inches(2.5), 0, Inches(2.5), Inches(0.4), BRICK)
add_rect(s, Inches(5.0), 0, Inches(2.5), Inches(0.4), MUSTARD)
add_rect(s, Inches(7.5), 0, Inches(2.5), Inches(0.4), MOSS)
add_rect(s, Inches(10.0), 0, Inches(3.33), Inches(0.4), PLUM)

# Color stripe bottom (mirror)
add_rect(s, 0, Inches(7.1), Inches(2.5), Inches(0.4), PLUM)
add_rect(s, Inches(2.5), Inches(7.1), Inches(2.5), Inches(0.4), MOSS)
add_rect(s, Inches(5.0), Inches(7.1), Inches(2.5), Inches(0.4), MUSTARD)
add_rect(s, Inches(7.5), Inches(7.1), Inches(2.5), Inches(0.4), BRICK)
add_rect(s, Inches(10.0), Inches(7.1), Inches(3.33), Inches(0.4), PETROL)

# End label
add_text(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.3),
         "END  ·  FIN  ·  THE END  ·  ありがとう",
         size=11, color=BRICK, bold=True, tracking=400, font=F_MONO)
add_line(s, Inches(0.8), Inches(1.35), Inches(2.5), Inches(1.35), color=BRICK, width=2)

# Massive Thank You
add_text(s, Inches(0.8), Inches(1.7), Inches(12), Inches(2.5),
         "Thank",
         size=160, bold=True, color=PETROL, font=F_DISPLAY)
add_text(s, Inches(0.8), Inches(3.5), Inches(12), Inches(2.0),
         "You.",
         size=160, bold=True, color=BRICK, font=F_DISPLAY)

# Subtitle
add_text(s, Inches(0.8), Inches(5.6), Inches(11), Inches(0.5),
         "感谢老师与同学们的聆听  ·  Thank you for your time and attention",
         size=15, color=INK_SOFT, italic=True)

# Live URL bar
url_y = Inches(6.0)
add_rect(s, Inches(0.8), url_y, Inches(12), Inches(0.55), INK)
add_text(s, Inches(1.0), url_y, Inches(2), Inches(0.55),
         "🌐 LIVE", size=9, color=MUSTARD, bold=True, tracking=300,
         anchor=MSO_ANCHOR.MIDDLE, font=F_MONO)
add_text(s, Inches(2.2), url_y, Inches(10.5), Inches(0.55),
         "signing-isle-printed-shapes.trycloudflare.com",
         size=12, color=WHITE, font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)

# GitHub repo bar
gh_y = Inches(6.6)
add_rect(s, Inches(0.8), gh_y, Inches(12), Inches(0.55), PETROL)
add_text(s, Inches(1.0), gh_y, Inches(2), Inches(0.55),
         "⌥ CODE", size=9, color=MUSTARD, bold=True, tracking=300,
         anchor=MSO_ANCHOR.MIDDLE, font=F_MONO)
add_text(s, Inches(2.2), gh_y, Inches(10.5), Inches(0.55),
         "github.com/a2318491287-design/campus-resource-platform",
         size=12, color=WHITE, font=F_MONO, anchor=MSO_ANCHOR.MIDDLE)

# Demo account bar
demo_y = Inches(6.7) - Inches(0.06)
add_text(s, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
         "演示账号  /  Demo Login:  学号 1230000000  ·  密码 demo123  ·  100 积分立即可用",
         size=11, color=INK_SOFT, italic=True)

prs.save('/Users/yuxianglian/Documents/系统分析与设计/SAD_Project/Final_Presentation.pptx')
print(f"Done: Final_Presentation.pptx — {len(prs.slides)} slides (editorial × multicolor)")
